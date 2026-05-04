from pathlib import Path
import sys
from pptx import Presentation

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.append(str(SCRIPT_DIR))

from ppt_theme_tech import style_title, style_subtitle, style_body, add_header_bar, add_highlight_box

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-monthly-ai-agent-summary', '2026-03-monthly-ai-agent-summary-v1.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = '2026-03 月 AI Agent 發展簡報'
prs.core_properties.author = 'ClawChan'


def add_slide(title, bullets, highlight=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    title_shape.text = title
    style_title(title_shape, 24)
    add_header_bar(slide)

    if highlight:
        add_highlight_box(slide, highlight, left=0.8, top=1.1, width=11.1, height=0.8)
        body_top = 2.1
    else:
        body_top = 1.4

    body = slide.placeholders[1]
    body.left = int(0.8 * 914400)
    body.width = int(11.1 * 914400)
    body.top = int(body_top * 914400)
    body.height = int(4.4 * 914400)

    tf = body.text_frame
    tf.clear()
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        first = False
    style_body(tf, 18)


# Cover
cover = prs.slides.add_slide(prs.slide_layouts[0])
cover.shapes.title.text = '2026-03 月 AI Agent 發展簡報'
cover.placeholders[1].text = '本月重點整理\nClawChan / 小爪\n2026-03-27'
style_title(cover.shapes.title, 28)
style_subtitle(cover.placeholders[1], 15)
add_header_bar(cover)

add_slide(
    '一句話總結',
    [
        '完成從安裝到可遠端操控、語音互動、流程化產出的閉環。',
        '把研究、報告、寄信、備份等工作變成可重複的流程。',
        '同時補強安全邊界與長任務穩定性。'
    ],
    highlight='本月關鍵：從「能用」提升到「可持續運作」。'
)

add_slide(
    '重要里程碑（3/13–3/21）',
    [
        '3/13：受啟發安裝 OpenClaw，專案啟動。',
        '3/14：多模型環境嘗試，學到配額與安裝差異。',
        '3/15：解決金鑰問題，成功完成首次提交。',
        '3/17：Telegram 遠端操控上線。',
        '3/18：子代理計時提醒成功。',
        '3/19：多代理辯論流程測試完成。',
        '3/20：切換到 OpenAI Codex。',
        '3/20：行動語音端到端 prototype 完成。',
        '3/21：新增語音控制密碼保護與寄信能力。'
    ]
)

add_slide(
    '能力增長：遠端 + 語音',
    [
        'Telegram 遠端操控穩定化。',
        '手機端語音 → 回覆 → TTS 端到端完成。',
        'cloudflared 方式保留外網存取彈性。',
        '語音長任務先回「收到/OK」避免斷線體感。'
    ],
    highlight='語音與遠端操控已可實際上線使用。'
)

add_slide(
    '流程化與多代理協作',
    [
        '研究流程：多角度、多輪辯論與交付串接。',
        '主代理需主動推進與回報進度。',
        '交付階段也需監控並可自動補救。',
        '固定化腳本與 quick-start 指南。'
    ],
    highlight='工作從「一次性」進化到「可重複流程」。'
)

add_slide(
    '工程化與穩定性補強',
    [
        '長任務改用背景模式，避免 SIGTERM。',
        '語音服務僅開放私人網段 + cloudflared。',
        '低風險命令加入 allowlist，降低打斷。',
        '成長日誌缺漏可於啟動時補檢。'
    ]
)

add_slide(
    '當前風險與下步',
    [
        '配額 / rate limit 仍是運作瓶頸。',
        'Approval UI 偶發失靈，需要更穩定的流程。',
        '後續：固定簡報版型、縮短交付回合時間。'
    ],
    highlight='重點：先把 approval 與配額穩定化。'
)

prs.save(OUT)
print(OUT)
