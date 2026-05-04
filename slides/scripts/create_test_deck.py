from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
import sys

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-22-powerpoint-workflow-test', '2026-03-22-powerpoint-workflow-test-v1.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = 'PowerPoint 工作模式測試'
prs.core_properties.subject = 'ClawChan PowerPoint workflow test'
prs.core_properties.author = 'ClawChan'

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = 'PowerPoint 工作模式測試'
slide.placeholders[1].text = 'ClawChan / 小爪\n2026-03-22'

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '目前環境已就緒'
text = slide.placeholders[1].text_frame
text.clear()
for i, line in enumerate([
    'LibreOffice Impress 已安裝',
    'python-pptx 已可用',
    '中文字型 Noto Sans CJK 已安裝',
    '之後可直接生成 .pptx 並寄送',
]):
    p = text.paragraphs[0] if i == 0 else text.add_paragraph()
    p.text = line
    p.level = 0

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '之後的標準流程'
text = slide.placeholders[1].text_frame
text.clear()
for i, line in enumerate([
    '先確認簡報目標、聽眾、時長',
    '先做大綱，再展開每頁內容',
    '生成 .pptx',
    '必要時補逐頁講稿與備註',
    '可再寄送 email 或繼續修訂',
]):
    p = text.paragraphs[0] if i == 0 else text.add_paragraph()
    p.text = line
    p.level = 0

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '備註'
box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(2.5))
frame = box.text_frame
p = frame.paragraphs[0]
p.text = '這是一份測試投影片，用來確認小爪已具備後續製作 PowerPoint 的能力。'
p.font.size = Pt(24)

prs.save(OUT)
print(OUT)
