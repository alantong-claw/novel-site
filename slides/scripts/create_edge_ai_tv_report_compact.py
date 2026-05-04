from pathlib import Path
from pptx import Presentation
import sys

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-22-edge-ai-tv-openclaw-report-compact-v1', '2026-03-22-edge-ai-tv-openclaw-report-compact-v1.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = '15 TOPS Edge AI + TV 控制 + OpenClaw 研究簡報（精簡版）'
prs.core_properties.author = 'ClawChan'


def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        first = False

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '15 TOPS Edge AI + TV 控制 + OpenClaw'
slide.placeholders[1].text = '研究簡報（精簡版）\nClawChan / 小爪\n2026-03-22'

add_slide('一句話結論', [
    '不要做萬用 AI 電視大腦。',
    '應做有邊界的 edge AV / control appliance。',
    'OpenClaw 適合當 orchestration layer，不是完整影音堆疊。'
])

add_slide('市場價值', [
    '價值不在 15 TOPS 本身，而在本地控制、低延遲、隱私與設備整合。',
    '真正能賣的是工作流程效率，不是單純算力。',
    '適合可控環境與可重複任務。'
])

add_slide('最有潛力的場景', [
    '會議室 / 教室',
    '數位看板 / kiosk',
    '安防 / 營運顯示牆',
    'Prosumer AV / 串流工作室',
    '無障礙顯示控制'
])

add_slide('技術主張', [
    '最佳情況：這台 edge box 自己就是播放來源。',
    '折衷方案：控制 TV，但不吃進所有影音內容。',
    '最差路徑：擷取任意 HDMI 來源，會踩 EDID / latency / HDCP 地雷。'
])

add_slide('OpenClaw 安裝建議', [
    '推薦 Linux 直接安裝。',
    '64-bit Linux + Node 22+/24 + openclaw onboard --install-daemon。',
    '再整合 CEC、播放、vendor control 等工具。'
])

add_slide('主要限制', [
    'DRM / HDCP 是硬邊界。',
    'CEC 可用但不穩，不能幻想通吃。',
    '消費級客廳 support 成本高。',
    '若過度賣 AI 幻想，容易掉進維運泥沼。'
])

add_slide('MVP 建議', [
    'Linux mini-PC / SBC + HDMI output + USB HDMI-CEC + mic/speaker。',
    '先做 local playback / local UI。',
    '先驗證：開關 display、切 source、語音控制 scene。'
])

prs.save(OUT)
print(OUT)
