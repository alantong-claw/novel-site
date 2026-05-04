from pathlib import Path
import sys
from pptx import Presentation

SCRIPT_DIR = Path(__file__).resolve().parent
if str(SCRIPT_DIR) not in sys.path:
    sys.path.append(str(SCRIPT_DIR))

from ppt_theme_tech import style_title, style_subtitle, style_body, add_header_bar, add_highlight_box

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-23-tv-edge-15tops-market', '2026-03-23-tv-edge-15tops-market-v3-tech.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = '15 TOPS TV Edge Chip 市場推廣研究'
prs.core_properties.author = 'ClawChan'


def add_slide(title, bullets, highlight=None):
    slide_layout = prs.slide_layouts[1] # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_shape = slide.shapes.title
    title_shape.text = title
    style_title(title_shape, 24) # Reduced font size for title
    
    add_header_bar(slide)

    # --- Content Placeholder Adjustments ---
    # Assume placeholder[1] is the content body
    body_placeholder = slide.placeholders[1]
    tf = body_placeholder.text_frame
    tf.clear()
    
    # Set placeholder dimensions for width and vertical alignment
    # These values might need fine-tuning based on default layout proportions
    # Standard margins are roughly 0.5 inch from each side. Placeholder width ~ 11.4 inches.
    body_placeholder.left = int(0.5 * 914400) # Approx 0.5 inch from left
    body_placeholder.width = int(11.4 * 914400) # Approx 11.4 inch width
    # body_placeholder.top = int(1.5 * 914400) # Default top adjusted if highlight exists

    body_top_default = 1.5
    body_top_with_highlight = 2.45

    if highlight:
        # Adjust highlight box and body top for better fit
        # Increased height and width for highlight box
        add_highlight_box(slide, highlight, left=0.7, top=1.35, width=11.4, height=1.8) # Increased height
        body_top = body_top_with_highlight
    else:
        body_top = body_top_default
    
    body_placeholder.top = int(body_top * 914400)

    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        first = False
    style_body(tf, 18)

# --- Cover Slide Adjustments ---
slide_cover = prs.slides.add_slide(prs.slide_layouts[0])

title_shape_cover = slide_cover.shapes.title
title_shape_cover.text = '15 TOPS 在 TV 上的 edge chip，能怎麼推廣市場？'

# Adjust title font size to prevent wrapping
style_title(title_shape_cover, 26) # Slightly reduced font size

# Ensure title placeholder has enough width
title_placeholder_cover = slide_cover.placeholders[0] # This is the title shape itself
# If title_shape_cover.left and title_shape_cover.width are accessible and modifiable, do it here. Often, the title shape is bound to the layout.

# For text in placeholder[1] (subtitle)
subtitle_placeholder = slide_cover.placeholders[1]
subtitle_placeholder.text = '多面向三輪研究後整理\nClawChan / 小爪\n2026-03-23'
style_subtitle(subtitle_placeholder, 15)

add_header_bar(slide_cover)

prs.save(OUT)
print(OUT)
