from pathlib import Path
from pptx import Presentation
import sys

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-22-edge-ai-tv-openclaw-report-compact-v2', '2026-03-22-edge-ai-tv-openclaw-report-compact-v2.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = '15 TOPS Edge AI + TV 控制 + OpenClaw 研究簡報（精簡版 v2）'
prs.core_properties.author = 'ClawChan'


def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        first = False

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '15 TOPS Edge AI + TV 控制 + OpenClaw'
slide.placeholders[1].text = '研究簡報（精簡版 v2）\nClawChan / 小爪\n2026-03-22'

add_slide('一句話結論', [
    '應做有邊界的 edge AV / control appliance。',
    '不要把它包裝成萬用 AI 電視大腦。',
    'OpenClaw 最適合做 orchestration layer。'
])

add_slide('價值在哪', [
    '本地控制、低延遲、隱私、離線能力。',
    '真正賣點是工作流程效率，不是單看 15 TOPS。',
    '適合可控環境與可重複任務。'
])

add_slide('最值得做的場景', [
    '會議室 / 教室',
    '數位看板 / kiosk',
    '安防 / 營運顯示牆',
    'Prosumer AV',
    '無障礙顯示控制'
])

add_slide('技術路線', [
    '最佳：edge box 自己就是播放來源。',
    '折衷：控制 TV，但不吃進所有影音。',
    '避免：任意 HDMI capture，容易踩 HDCP / EDID / latency 地雷。'
])

add_slide('OpenClaw 安裝與限制', [
    '推薦 Linux 直接安裝。',
    '本體不重，但周邊整合才是真難點。',
    'CEC、capture、vendor control 都需要外部工具補上。'
])

add_slide('MVP 建議', [
    'Linux mini-PC / SBC + HDMI output + USB HDMI-CEC。',
    'OpenClaw 當 voice / automation / orchestration layer。',
    '先驗證：開關 display、切 source、語音控制 scene。'
])

prs.save(OUT)
print(OUT)
