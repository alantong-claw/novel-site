from pathlib import Path
from pptx import Presentation
import sys

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-23-digital-ic-ai-adoption', '2026-03-23-digital-ic-ai-adoption-v1.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = '如何推廣數位 IC 設計工程師使用 AI 輔助工作'
prs.core_properties.author = 'ClawChan'

def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        first = False

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '如何推廣數位 IC 設計工程師使用 AI 輔助工作'
slide.placeholders[1].text = '三輪正反研究後整理\nClawChan / 小爪\n2026-03-23'

add_slide('一句話結論', [
    '可以推，但不能從「AI 幫你設計晶片」切入。',
    '應從低風險、可審查、可量測的周邊工程工作切入。',
    '先證明可信 ROI，再逐步擴張。'
])
add_slide('為什麼數位 IC 團隊不會自然採用 AI', [
    '正確性風險不對稱：錯誤代價高。',
    '工作高度依賴專案上下文。',
    '現有流程已高度優化。',
    'IP / security 是硬門檻。',
    '前端快，不代表 end-to-end throughput 真變快。'
])
add_slide('最值得先推的場景', [
    'Regression / log triage',
    'Spec / internal knowledge retrieval',
    'Script / automation drafting',
    'Documentation / review support',
    'Testbench / assertion scaffolding（需人審）'
])
add_slide('不建議早期推的場景', [
    'Critical RTL generation',
    'SDC / STA authoring',
    'CDC / RDC waiver reasoning',
    'ECO proposal near tapeout',
    'Signoff judgment / DFT / UPF 關鍵決策'
])
add_slide('建議的 rollout 方式', [
    'Phase 0: guardrails / security / allowed use cases',
    'Phase 1: 個人生產力 pilot',
    'Phase 2: engineer-in-the-loop artifact pilot',
    'Phase 3: workflow integration',
    'Phase 4: controlled scale-up'
])
add_slide('應量測的指標', [
    'time to first useful answer',
    'regression triage time',
    'script / testbench first draft time',
    'review rejection / rework / escaped defect',
    'weekly active users / repeat use',
    'senior interruption load / security incidents'
])
add_slide('Go / No-Go', [
    '可擴張：10–20% task-level 改善、品質不退步、repeat use 穩定。',
    '應暫停：defect / rework 上升、review burden 上升、security 邊界不清。',
    '重點不是 AI 有沒有生成內容，而是有沒有提升 trusted engineering throughput。'
])
add_slide('最後建議', [
    '把 AI 當成工程生產力增幅層，而不是設計 authority。',
    '先拿掉高摩擦、低風險的重複工作。',
    '先建立可信度，再談擴張。'
])

prs.save(OUT)
print(OUT)
