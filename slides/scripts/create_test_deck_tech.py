from pathlib import Path
from pptx import Presentation
from ppt_theme_tech import style_title, style_subtitle, style_body, add_header_bar, add_highlight_box
import sys

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('tech-style-demo', 'tech-style-demo.pptx', 'OUTPUT_PPT')

prs = Presentation()

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '簡潔科技風 Demo'
slide.placeholders[1].text = 'ClawChan slide theme\n2026-03-24'
style_title(slide.shapes.title, 28)
style_subtitle(slide.placeholders[1], 15)
add_header_bar(slide)

slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = '關鍵結論'
style_title(slide2.shapes.title, 24)
add_header_bar(slide2)
add_highlight_box(slide2, '這套風格主打：簡潔、科技感、可重用，不靠花俏圖片也能提升層次。')
body = slide2.placeholders[1].text_frame
body.text = '• 主色深藍\n• 重點用淺藍框\n• 適合研究、提案、比較型簡報'
style_body(body, 18)

prs.save(OUT)
print(OUT)
