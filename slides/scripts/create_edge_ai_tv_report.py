from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path('/home/alantong/ai-work/slides')
DEFAULT_OUT = Path('/home/alantong/ai-work/work_tmp/tasks/2026-03-22-edge-ai-tv-openclaw-report/2026-03-22-edge-ai-tv-openclaw-report-v1.pptx')
OUT = Path(__import__('os').environ.get('OUTPUT_PPT', str(DEFAULT_OUT)))
OUT.parent.mkdir(parents=True, exist_ok=True)

prs = Presentation()
prs.core_properties.title = '15 TOPS Edge AI + TV 控制 + OpenClaw 研究整理'
prs.core_properties.subject = 'Edge AI TV control research summary'
prs.core_properties.author = 'ClawChan'


def add_bullets_slide(title, bullets, level1=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = text
        p.level = level
        first = False
    return slide

# 1 cover
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '15 TOPS Edge AI + TV 控制 + OpenClaw 研究整理'
slide.placeholders[1].text = '根據 2026-03-21 研究結果整理\nClawChan / 小爪\n2026-03-22'

# 2 thesis
add_bullets_slide('一句話結論', [
    '真正有價值的，不是「萬用 AI 電視大腦」，而是「有邊界的 edge AV / control appliance」。',
    '賣點不是 15 TOPS 本身，而是本地控制、低延遲、隱私、可離線與硬體整合。',
    'OpenClaw 最適合當 orchestration / agent layer，不應被誤當成完整影音堆疊。'
])

# 3 market value
add_bullets_slide('市場價值在哪', [
    '價值來自：本地推論 + TV / display 控制 + 語音互動 + 場景感知 + 實體設備整合',
    '不是單賣算力，而是替真實場景省操作成本、切換成本、反應時間與維運成本',
    '最適合落在「可控環境、可重複流程、可驗證硬體矩陣」的應用場景'
])

# 4 applications
add_bullets_slide('最有潛力的應用場景', [
    '會議室 / 教室：切換來源、啟停顯示、錄影 / 串流工作流',
    '數位看板 / kiosk / 零售螢幕：來源切換、fallback content、本地分析',
    '安防 / 營運儀表板：事件觸發顯示切換、本地 triage',
    'Prosumer AV / 串流工作室：scene / source switching、語音工作流',
    '無障礙 / 輔助控制：簡化 TV / display / volume / source 操作'
])

# 5 access modes
add_bullets_slide('存取電視影音：三種層次', [
    '最佳：裝置自己就是播放來源（最穩、最容易理解與控制）',
    '折衷：控制 TV，但不真的吃進所有影音內容（room-level copilot）',
    '最難：擷取任意 HDMI / TV 來源分析（capture / EDID / latency / HDCP 問題多）'
])

# 6 cec
add_bullets_slide('控制電視：HDMI-CEC 是主路徑，但不能迷信', [
    'CEC 可做到：開關機、切輸入源、音量控制、部分遙控訊號',
    'Linux 常見方案：kernel CEC、libCEC、USB-CEC adapter',
    '但不同 TV / soundbar / AVR 差異大，必須當成「驗證過硬體矩陣上的功能」',
    '產品化應準備 fallback：vendor API、IP control、Wake-on-LAN、HID / serial / relay'
])

# 7 15 tops
add_bullets_slide('15 TOPS 真正適合做什麼', [
    'wake word / trigger detection',
    '輕量 ASR 輔助、OCR、scene / event detection',
    'speaker / activity detection、本地 routing / triage',
    '不應過度期待：大模型等級多模態理解、只看 TOPS 就保證好體驗'
])

# 8 openclaw install
add_bullets_slide('OpenClaw（龍蝦）怎麼裝', [
    '推薦：Linux 直接安裝',
    ('64-bit Linux + Node 22+/24', 1),
    ('npm i -g openclaw@latest', 1),
    ('openclaw onboard --install-daemon', 1),
    '再整合 CEC、播放、capture、vendor control 工具',
    '優點：最好接硬體、最好 debug、最適合做成 appliance-like service',
    '缺點：OpenClaw 不是完整影音中介層，穩定度高度依賴周邊 stack'
])

# 9 deployment
add_bullets_slide('部署模式建議', [
    'A. fully local edge-first：低延遲、硬體整合最好，但維運壓力較大',
    'B. central OpenClaw + edge box 只處理 I/O：集中管理較好，但延遲與連線依賴較高',
    '建議實務：edge-first 做本地 I/O，central assist 做管理與 orchestration'
])

# 10 limitations
add_bullets_slide('主要限制與盲點', [
    'DRM / HDCP 是硬邊界',
    'CEC 相容性不穩，不能幻想通吃',
    '消費級客廳 support 成本很高',
    '若賣點只是「15 TOPS 很強」，很容易 oversell',
    'OpenClaw 帶來控制力，也帶來邊界管理責任'
])

# 11 mvp
add_bullets_slide('建議 MVP', [
    '不要做萬用 AI 電視盒',
    '要做有邊界的 AV control appliance',
    'Linux mini-PC / SBC + HDMI output + USB HDMI-CEC + mic/speaker',
    'OpenClaw 當 orchestration / voice / automation layer',
    'MVP 先從 local playback / local UI 做起',
    '先驗證：開/關 display、切 source、語音控制 scene、事件切換顯示內容'
])

# 12 go no-go
add_bullets_slide('Go / No-Go 判準', [
    'Go：部署環境可控、硬體矩陣窄、價值來自本地控制/隱私/低延遲、流程可重複',
    'No-Go：想通吃所有 TV / HDMI / 影音來源，或依賴大規模讀受保護內容',
    '第一站不建議打 unmanaged consumer 市場'
])

# 13 closing
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '最後一句'
box = slide.shapes.add_textbox(Inches(0.9), Inches(1.8), Inches(8.2), Inches(2.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = '做對了，它會是有明確場景價值的 edge AV/control 產品；做錯了，就會掉進 HDMI、DRM、CEC 與 support 成本的泥沼。'
p.font.size = Pt(24)

prs.save(OUT)
print(OUT)
