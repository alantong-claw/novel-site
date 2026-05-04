from pathlib import Path
from pptx import Presentation
import sys

ROOT = Path('/home/alantong/ai-work/slides')
HELPER_DIR = Path('/home/alantong/ai-work/scripts')
if str(HELPER_DIR) not in sys.path:
    sys.path.append(str(HELPER_DIR))
from office_output_path import output_path
OUT = output_path('2026-03-22-hsinchu-science-vs-gifted', '2026-03-22-hsinchu-science-vs-gifted-v1.pptx', 'OUTPUT_PPT')

prs = Presentation()
prs.core_properties.title = '新竹實驗高中科學班 vs 資優班'
prs.core_properties.author = 'ClawChan'


def add_slide(title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    first = True
    for item in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        p.text = item
        p.level = 0
        first = False

slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '新竹實驗高中科學班 vs 資優班'
slide.placeholders[1].text = '三輪正反辯論後整理\nClawChan / 小爪\n2026-03-22'

add_slide('一句話結論', [
    '這不是哪一班比較強，而是哪個環境更適合這個學生。',
    '很確定走 STEM、喜歡集中火力的人：科學班較有利。',
    '很強但仍在探索、想保留彈性的人：資優班較穩健。'
])

add_slide('科學班的主要優勢', [
    '數理訓練更集中、更深，對 STEM 明確導向者更有效率。',
    '同儕更偏數理、研究、競賽，學習氛圍集中。',
    '較容易接到科展、研究、競賽等資源。',
    '若未來已偏醫學、工程、科學研究，提早專精有利。'
])

add_slide('科學班的主要風險', [
    '壓力較集中，容易 burnout。',
    '發展可能較窄。',
    '若其實還在探索期，過早鎖定風險較高。'
])

add_slide('資優班的主要優勢', [
    '彈性高，保留更多未來選擇。',
    '更容易兼顧理工、表達、領導、跨域活動。',
    '壓力通常較分散，不像單一路徑淘汰賽。',
    '對尚未完全確定方向的學生更友善。'
])

add_slide('資優班的主要風險', [
    '若學生已非常明確偏 STEM，可能覺得不夠集中。',
    '研究與競賽資源更依賴個人主動爭取。',
    '若自律不足，容易浪費彈性優勢。'
])

add_slide('真正的分水嶺', [
    '是否已明確偏向 STEM？',
    '是否喜歡高密度數理環境？',
    '是否能承受較強比較與壓力？',
    '還是更需要探索空間與整體發展？'
])

add_slide('建議怎麼選', [
    '選科學班：已明確偏 STEM、喜歡深挖、想走研究/競賽/醫工理科。',
    '選資優班：很強但方向未定、想保留彈性、重視全面發展。',
    '關鍵不是名聲，而是適配度。'
])

prs.save(OUT)
print(OUT)
